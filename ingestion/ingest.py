#<------------------------------------HEADER------------------------------------------->
#Imports
import time
from datetime import datetime
from pathlib import Path
import os
import zipfile
import tempfile
import json
import boto3
from modal import App, Image, Secret
from pinecone import Pinecone, Index, ServerlessSpec
import openpyxl
import extract_msg
import requests
import base64
import pypdf
from bs4 import BeautifulSoup
from openai import OpenAI
import xlrd
import pdfminer
import nltk
import csv
from pdfminer.high_level import extract_text
from pptx import Presentation
from docx import Document
from langchain_openai import OpenAIEmbeddings
from dotenv import load_dotenv
load_dotenv()
#Creating an instance of OpenAI class
client = OpenAI()
nltk.download('stopwords')
# Setup the Modal image with necessary Python packages and environment variables
image = Image.debian_slim().pip_install(
    "pypdf==4.2.0",
    "pinecone-client==4.0.0",
    "langchain-openai==0.1.6",
    "langchain==0.1.12",
    "langchain-community==0.0.36",
    "boto3==1.20.21",
    "unstructured==0.11.8",
    "unstructured-client==0.22.0",
    "requests>=2.25.1",
    "python-pptx",
    "python-docx",
    "openpyxl",
    "beautifulsoup4",
    "extract_msg",
    "openai",
    "bs4",
    "xlrd",
    "python-dotenv",
    "pdfminer.six",
    "nltk",
    "python-pptx",
    "python-docx",
).env({
    "PINECONE_API_KEY": os.getenv("PINECONE_API_KEY"),
    "OPENAI_API_KEY": os.getenv("OPENAI_API_KEY"),
    "AWS_ACCESS_KEY_ID": os.getenv("AWS_ACCESS_KEY_ID"),
    "AWS_SECRET_ACCESS_KEY": os.getenv("AWS_SECRET_ACCESS_KEY"),
    "AWS_DEFAULT_REGION": "us-east-1"
})        
#Creating a Modal app
app = App(name="pinecone-ingestor", image=image)
#<------------------------------------BODY------------------------------------------->
class LangchainPineconeLoader:
    """
    LangchainPineconeLoader class is defined, which includes methods
    for processing different file types and interacting with S3 and Pinecone.

    Attributes:
        bucket_name (str): Name of the S3 bucket from which files are loaded.
        directory_path (str): The directory path in the S3 bucket to start file retrieval.
        index_name (str): Name of the Pinecone index for storing vector embeddings.
        batch_size (int): Number of documents to process in each batch.
        embedding_model (str): Model identifier for generating text embeddings using OpenAI.

    Methods:
        load_and_index(): Loads and processes files from the specified S3 bucket and indexes their content.
        process_zip_file(zip_path): Extracts and processes files from a zip archive.
        _load_and_split_file(file_path, ext): Loads a file based on its extension and processes its content.
        generate_embedding(text): Generates vector embeddings for a given text.
        custom_upsert(data): Performs a custom upsert operation into Pinecone with generated embeddings.
        process_excel_file(file_path): Processes Excel files (.xls, .xlsx) to extract text content.
        process_image_file(image_path): Processes image files to extract textual descriptions.
        process_msg_file(file_path): Processes .msg files to extract their text content.
        index_texts_into_pinecone(texts): Inserts batches of text into the Pinecone index.
    """
    def __init__(self, bucket_name, directory_path, index_name="default-index", batch_size=100, embedding_model='text-embedding-3-small'):
        """
        Initializes the loader with specified parameters and creates a Pinecone index if it does not exist.

        Args:
            bucket_name (str): Name of the S3 bucket from which files are loaded.
            directory_path (str): The directory path in the S3 bucket to start file retrieval.
            index_name (str): Name of the Pinecone index for storing vector embeddings. Defaults to "default-index".
            batch_size (int): Number of documents to process in each batch. Defaults to 100.
            embedding_model (str): Model identifier for generating text embeddings using OpenAI. Defaults to 'text-embedding-3-small'.
        """
        self.bucket_name = bucket_name
        self.directory_path = directory_path.strip('/')
        self.index_name = index_name
        self.batch_size = batch_size
        self.pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))
        self.embeddings = OpenAIEmbeddings(model=embedding_model)
        self.total_sheets_processed=0
        self.total_vectors=0
        self.files_processed=0
        self.index_id =0
        self.filename = "default"
        self.chunk_size = 8000
        self.loaders = {
            '.pdf': self.process_pdf_file,
            '.csv': self.process_csv_file,
            '.txt': self.process_text_file,
            '.json': self.process_json_file,
            '.docx': self.process_docx_file,
            '.pptx': self.process_pptx_file,
            '.xlsx': self.process_excel_file,
            '.xls': self.process_excel_file,
            '.png': self.process_image_file,
            '.msg': self.process_msg_file,
        }
        index_names = [index['name'] for index in self.pc.list_indexes()]
        if self.index_name not in index_names:
            self.pc.create_index(
                name=self.index_name,
                dimension=1536,
                metric="euclidean",
                spec=ServerlessSpec(cloud="aws", region="us-east-1")
            )
            print("Successfully created the new index with name", self.index_name)

#Loading files from S3 
    def load_and_index(self):
        """
        Main method to load files from S3, process them, and index their contents.
        Handles the file discovery, downloading, and processing logic for various file types.
        It iterates over the objects in the specified S3 directory, downloads each file,
        determines the appropriate processing method based on the file extension, and indexes the extracted content.
        """
        s3_client = boto3.client('s3')
        paginator = s3_client.get_paginator('list_objects_v2')
        page_iterator = paginator.paginate(Bucket=self.bucket_name, Prefix=self.directory_path) 
        
        found_objects = False
        for page in page_iterator:
            if 'Contents' in page:  # Check if 'Contents' key exists
                for obj in page['Contents']:
                    file_key = obj['Key']
                    self.filename=file_key
                    if not file_key.endswith('/'):  # Correct method used here
                        _, ext = os.path.splitext(file_key)  # Extract extension from the original file name
                        print(file_key)
                        # Create a temporary file with the original extension
                        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_file:
                            s3_client.download_file(self.bucket_name, file_key, temp_file.name)
                            if file_key.lower().endswith('.zip'):  # Correct method used here
                                self.process_zip_file(temp_file.name)
                            else:
                                print(temp_file.name) 
                                self._load_and_split_file(temp_file.name, ext,file_key) 
                                print(f"The file {file_key} has been successfully upserted to pinecone")
                                print("----------------------------------0000000000----------------------0000000000000000--------------")
                                self.files_processed+=1
                                 # Pass extension explicitly if needed
                            os.unlink(temp_file.name)  # Clean up the temporary file
                found_objects = True
            else:
                print("No files found at the specified path.")

        if not found_objects:
            print("No contents available in the bucket/prefix.")
        self.finalize_processing()

    def finalize_processing(self):
        print(f"Total number of sheets processed: {self.total_sheets_processed}")
        print(f"total vectors upserted == {self.total_vectors}")
        print(f"total files processed == {self.files_processed}")

#Processing and extracting text from different file formats and specific methods like process_pdf_file, process_excel_file.
    def _load_and_split_file(self, file_path,ext=None,filename=None):
        """
        Determines the appropriate loader based on the file extension and processes the file to extract text.

        Args:
            file_path (str): The full path to the file to be processed.
            ext (str, optional): The file extension. If not provided, it is derived from the file path.

        This method uses a mapping of file extensions to loader classes to process the file
        and extract text based on the file type. It handles different types of documents and
        special cases like images and spreadsheets differently.
        """
        ext = ext if ext else os.path.splitext(file_path)[1]
        ext=ext.lower()
        print(f"Processing file: {file_path} with extension: {ext}")
        loader_cls = self.loaders.get(ext)
        if loader_cls:
            #for msg files
            if loader_cls==self.process_msg_file:
                #passing both file_path and filename
                texts=self.process_msg_file(file_path,filename)
                self.custom_upsert(texts,filepath=filename) 
                print(f"Successfully upserted {filename}")
                #for img files
            elif loader_cls== self.process_image_file:
                    print("considering image file")
                    #passing both file_path and filename
                    self.custom_upsert(loader_cls(file_path,filename),filepath=filename)
                    print(f"Successfully upserted {filename}")
            #for excel filess
            elif loader_cls==self.process_excel_file: 
                #passing both file_path and filename
                loader_cls(file_path,filename)
                print(f"Successfully upserted {filename}") # Custom method to handle upserting differently
            elif loader_cls==self.process_pptx_file:
                loader_cls(file_path)
            elif loader_cls==self.process_docx_file:
                loader_cls(file_path)
            elif loader_cls==self.process_csv_file:
                loader_cls(file_path)
            elif loader_cls==self.process_text_file:
                loader_cls(file_path)
            elif loader_cls==self.process_json_file:
                loader_cls(file_path)
            elif loader_cls==self.process_pdf_file:
                try:
                    loader_cls(file_path)
                except Exception as e:
                    print(f"STATUS CODE:{e}")
                    print(f"Error loading file {file_path}: {e}")
            else:
                print(f"No loader found for extension {ext}")
        else:
            print(f"No loader found for extension {ext}")
        return []
    

    def process_zip_file(self, zip_path):
        """
        Processes each file within a zip archive by extracting its contents and processing each file.
        Args:
            zip_path (str): The file path of the zip archive to be processed.
        This method extracts all files in the zip archive to a temporary directory, processes each file
        for text extraction, and cleans up the extracted files after processing.
        """
        with zipfile.ZipFile(zip_path, 'r') as z:
            z.extractall(path=os.path.dirname(zip_path))
            for filename in z.namelist():
                file_path = os.path.join(os.path.dirname(zip_path), filename)
                if not os.path.isdir(file_path):
                    texts = self._load_and_split_file(file_path)
                    if texts:
                        self.index_texts_into_pinecone(texts)
                    os.unlink(file_path)  # Clean up the extracted files

    def split_text_into_chunks(self,text, chunk_size=2000):
        return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]

    def process_excel_file(self, file_path,filename = None):
        """
        Processes Excel files, extracting text content from cells and compiling it into a single string.
        Args:
            file_path (str): The file path of the Excel file to be processed.
        Returns:
            str: The concatenated text content of all cells in the Excel file.

        This method handles both .xlsx and .xls file formats. Openpyxl for .xlsx and xlrd for .xls are used to open and read sheets
        and rows. It compiles all cell content into a single string which is then used for further processing.
        """
        extension = os.path.splitext(file_path)[1]
        if extension == ".xlsx":
            workbook = openpyxl.load_workbook(file_path, read_only=True)
            for sheet_name in workbook.sheetnames:
                lines = []
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows(values_only=True):
                    row_content = [str(cell) if cell is not None else "" for cell in row]
                    lines.append(' '.join(row_content))
                text_content = '\n'.join(lines)
                document_id = filename
                self.custom_upsert({'id': document_id, 'content': text_content},sheetname=sheet_name,filepath=filename)

        elif extension == ".xls":
            workbook = xlrd.open_workbook(file_path, on_demand=True)
            for sheet_name in workbook.sheet_names():
                lines = []
                sheet = workbook.sheet_by_name(sheet_name)
                for row_idx in range(sheet.nrows):
                    row = sheet.row(row_idx)
                    row_content = [str(cell.value) if cell.value is not None else "" for cell in row]
                    lines.append(' '.join(row_content))
                text_content = '\n'.join(lines)
                document_id = filename
                self.custom_upsert({'id': document_id, 'content': text_content},sheetname=sheet_name,filepath=filename)

 
    def process_image_file(self,image_path,filename=None):
        """
        Processes image files to generate a text description.
        Args:
            image_path (str): The file path of the image to be processed.
        Returns:
            str: A detailed description of the image content.
        This method reads the image file, converts it to a base64 string, and sends it to GPT4Vision
        to obtain a text description of the image contents.
        """
        print("PROCESS_IMAGE FUNCTION CALLED")
        with open(image_path, "rb") as image_file:
            print("BASE64 CONVERSION DONE")
            base64_image = base64.b64encode(image_file.read()).decode("utf-8")
            #print(base64_image)
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {os.getenv('OPENAI_API_KEY')}"
        }
        payload = {
            "model": "gpt-4-vision-preview",
            "messages": [{
                "role": "user",
                "content": [{
                    "type": "text",
                    "text": "return the image content, only the content and not a single word other than the image content in the response.Please be as much detailed and comprehensive as possible"
                }, {
                    "type": "image_url",
                    "image_url": f"data:image/jpeg;base64,{base64_image}"
                }]
            }],
            "max_tokens": 3000
        }
        print("RESPONSE REQUEST SENT")
        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
        print("RESPONSE RECIEVED")
        #print(f"The response from OPENAI IS ---------------------------->",response)
        if response.status_code == 200:
            content= response.json()['choices'][0]['message']['content']
            return {'id': filename, 'content': content}
        else:
            return "Failed to get image description from GPT-4"


    def process_msg_file(self, file_path,filename=None):
        """
        Processes .msg files to extract and return their textual content.
        Args:
            file_path (str): The file path of the .msg file to be processed.
        Returns:
            str: The extracted textual content from the .msg file.
        This method uses a BeautifulSoup4 to open and extract both plain text and HTML content from .msg files.
        """
        msg = extract_msg.Message(file_path)
        html_body_bytes = msg.htmlBody
        if not html_body_bytes:  # Check if the HTML body is empty
            print(f"Skipping processing of empty MSG file: {file_path}")
            return "No content in MSG file"
        html_body = html_body_bytes.decode('utf-8', errors='ignore')  # Decode bytes to string, ignoring errors
        soup = BeautifulSoup(html_body, 'html.parser')  # Parse HTML content directly using BeautifulSoup
        text_content = ' '.join(soup.stripped_strings)  # Concatenate all text into a single string
        if text_content:
            return {'id': filename, 'content': text_content}
        else:
            return None
        
    def preprocess_text(self, text):
        stop_words = set(nltk.corpus.stopwords.words('english'))
        words = [w.lower() for w in text.split() if w not in stop_words]
        return " ".join(words)

    def process_pdf_file(self,file_path):
        text=self.preprocess_text(extract_text(file_path))
        self.custom_upsert({'id':0, 'content': text},filepath=self.filename)

    def process_pptx_file(self,pptx_path):
        text_content = ""
        try:
            prs = Presentation(pptx_path)  # Load the presentation directly from the file path
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
            self.custom_upsert({'id':0, 'content': text_content},filepath=self.filename)
        except Exception as e:
            print(f"Error processing PowerPoint file '{pptx_path}': {e}")

    def process_docx_file(self,docx_path):
        try:
            text_content = []
            doc = Document(docx_path)  # Load the document from the file path
            for para in doc.paragraphs:
                print("Before the insertion",text_content)
                print("The actual content is",para.text)
                text_content.append(para.text)  # Append paragraph text to the list
                print("After the insertion",text_content)
            text_content= '\n'.join(text_content)  # Join all paragraphs to form a single string
            self.custom_upsert({'id':0, 'content': text_content},filepath=self.filename)
        except Exception as e:
            print(f"Error processing Word file '{docx_path}': {e}")

    def process_csv_file(self, csv_path, batch_size=100):
        try:
            concatenated_row=''
            with open(csv_path, mode='r', newline='', encoding='utf-8') as file:
                reader = csv.reader(file)
                for row in reader:
                    concatenated_row += ' '.join(row)+' '    
            self.custom_upsert({'id':0, 'content': concatenated_row},filepath=self.filename)

        except Exception as e:
            print(f"Error processing CSV file '{csv_path}': {e}")
    
    def process_text_file(self,text_path):
        try:
            with open(text_path, 'r', encoding='utf-8') as file:
                text = file.read()
                self.custom_upsert({'id':0, 'content': text},filepath=self.filename)
        except Exception as e:
            print(f"Error processing text file '{text_path}': {e}")

    def process_json_file(self,json_path):
        try:
            with open(json_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                json_string = json.dumps(data)
                self.custom_upsert({'id':0, 'content': json_string},filepath=self.filename)
        except Exception as e:
            print(f"Error processing JSON file '{json_path}': {e}")

#Generating text embeddings using OpenAI's API 
    def generate_embedding(self, text):
        """
        Generates and returns a vector embedding for the given text using the specified embedding model.

        Args:
            text (str): The text to be embedded.

        Returns:
            list: A list containing the numerical vector embedding of the input text.

        This method interfaces with OpenAI's API to convert text into a high-dimensional vector
        representing its semantic meaning.
        """
        response = client.embeddings.create(input=text, model='text-embedding-3-small')
        return response.data[0].embedding

#Upserting the generated embeddings into Pinecone
    def custom_upsert(self, document,filepath=None,sheetname=None):
        try:
            if sheetname:
                print(f"THIS IS A EXCEL FILE SHEET, SHEET NAME = {sheetname}")
                print("####################################################")
            # Manually split the document content into chunks
            text_chunks = self.split_text_into_chunks(document['content'], chunk_size=4000)
            index = self.pc.Index(self.index_name)
            upsert_count = 0
            # Determine the filename with extension based on sheetname and filepath
            if sheetname:
                filename_with_extension = f"{os.path.splitext(os.path.basename(filepath))[0]}_{sheetname}"
            else:
                filename_with_extension = os.path.basename(filepath) if filepath else document['id']

            # Generate and upsert embeddings for each chunk
            for chunk in text_chunks:
                embedding = self.generate_embedding(chunk)
                upsert_count += 1
                vector_id=f'{self.index_id+1}'
                self.index_id+=1
                if sheetname:
                    metadata = {'Chunk_ID': upsert_count,'sheetname': sheetname, 'filename': filepath,'content':chunk}
                else:
                    metadata={'Chunk_ID':upsert_count,'filename':filepath,'content':chunk}
                index.upsert(vectors=[(vector_id, embedding,metadata)])
                print(f"Upserted chunk {upsert_count} of document {vector_id} into Pinecone.")

            print(f"Total {upsert_count} parts upserted for document {filename_with_extension}.")
            # Increment the total sheet count if this upsert involves an Excel sheet
            if sheetname:
                self.total_sheets_processed += 1
            self.total_vectors+=len(text_chunks)
        except Exception as e:
            print(f"Error upserting document {document['id']}: {e}")
#<------------------------------------FOOTER------------------------------------------->
#loader_func is defined as a Modal function that initializes and runs the LangchainPineconeLoader.
@app.function(timeout=4000)
def loader_func():
    """
    Defines a loader function that initializes and runs the LangchainPineconeLoader.

    This function is decorated to specify a timeout, meaning the function must complete its execution
    within 8000 milliseconds. It initializes the loader with specific parameters for the S3 bucket,
    directory path, and Pinecone index name, and then starts the loading and indexing process.

    The parameters are:
    - bucket_name: 'transilience' specifies the S3 bucket from which files are loaded.
    - directory_path: 'compliance' specifies the path within the bucket where files are located.
    - index_name: 'final-test-8' specifies the name of the Pinecone index where vectors will be stored.
    """
    loader = LangchainPineconeLoader(bucket_name="ncrbtesting", directory_path="FINAL NCRB DATA /", index_name="testing-ncrb")
    loader.load_and_index()

#run_loader serves as the entry point for running the loader_func.
@app.local_entrypoint()
def run_loader():
    """
    Defines entry point for the application.
    This function acts as the entry point for running the `loader_func`.
    """
    loader_func.remote()
if __name__ == "__main__":
    run_loader()